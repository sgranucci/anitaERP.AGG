#!/usr/bin/env python3
"""Genera PowerPoint: Anita ERP — Módulo Gastronomía (desarrollo completo)."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = "/var/www/html/anitaERP/docs/presentaciones/AnitaERP_Gastronomia_Desarrollo_Completo.pptx"
IMG = Path("/var/www/html/anitaERP/docs/presentaciones/img_informe")

NAVY = RGBColor(0x1F, 0x4E, 0x79)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
ACCENT = RGBColor(0xC0, 0x59, 0x11)
LIGHT = RGBColor(0xF2, 0xF6, 0xFA)
DARK = RGBColor(0x17, 0x20, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5D, 0x6D, 0x7E)
GREEN = RGBColor(0x1E, 0x84, 0x4E)
TEAL = RGBColor(0x0E, 0x66, 0x5A)
ORANGE = RGBColor(0xD3, 0x54, 0x00)
RED = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 32


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def image_slide(title, image_path, page, caption=None, max_h=5.5):
    """Slide with title bar + centered image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, 22, True, WHITE)

    path = Path(image_path)
    if path.exists():
        # fit image under header
        top = Inches(1.15)
        left = Inches(0.5)
        max_w = Inches(12.3)
        max_h_emu = Inches(max_h)
        pic = slide.shapes.add_picture(str(path), left, top)
        # scale to fit
        if pic.width > max_w:
            ratio = max_w / pic.width
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
        if pic.height > max_h_emu:
            ratio = max_h_emu / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
        pic.left = int((prs.slide_width - pic.width) / 2)
        pic.top = int(top + (max_h_emu - pic.height) / 2) if pic.height < max_h_emu else int(top)
    if caption:
        ct = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.35))
        cr = ct.text_frame.paragraphs[0].add_run()
        cr.text = caption
        set_run(cr, 11, False, GRAY)
    add_footer(slide, page)


def two_image_slide(title, left_img, right_img, page, left_cap="", right_cap=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, 22, True, WHITE)

    slots = [
        (Inches(0.3), left_img, left_cap),
        (Inches(6.85), right_img, right_cap),
    ]
    for left, img, cap in slots:
        path = Path(img)
        box_w = Inches(6.15)
        box_h = Inches(5.35)
        top = Inches(1.2)
        if path.exists():
            pic = slide.shapes.add_picture(str(path), left, top)
            if pic.width > box_w:
                ratio = box_w / pic.width
                pic.width = int(pic.width * ratio)
                pic.height = int(pic.height * ratio)
            if pic.height > box_h:
                ratio = box_h / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = int(pic.height * ratio)
            pic.left = int(left + (box_w - pic.width) / 2)
            pic.top = int(top + (box_h - pic.height) / 2)
        if cap:
            ct = slide.shapes.add_textbox(left, Inches(6.6), box_w, Inches(0.35))
            cp = ct.text_frame.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            cr.text = cap
            set_run(cr, 11, True, BLUE)
    add_footer(slide, page)


def add_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_bar(slide, top=0, height=Inches(1.0), color=NAVY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, top, prs.slide_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, page):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.15), prs.slide_width, Inches(0.35)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(10), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Anita ERP  ·  Módulo Gastronomía  ·  Desarrollo completo"
    set_run(run, 10, False, WHITE)
    tb2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.18), Inches(1.5), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = f"{page} / {TOTAL}"
    set_run(run2, 10, False, WHITE)


def title_slide(title, subtitle, date_line):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.4), Inches(0.25), Inches(2.6)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, 40, True, WHITE)
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.9))
    run2 = tb2.text_frame.paragraphs[0].add_run()
    run2.text = subtitle
    set_run(run2, 20, False, RGBColor(0xD6, 0xE3, 0xF0))
    tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5))
    run3 = tb3.text_frame.paragraphs[0].add_run()
    run3.text = date_line
    set_run(run3, 14, False, RGBColor(0xA8, 0xC0, 0xD8))


def section_slide(number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    tb0 = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11), Inches(0.5))
    run0 = tb0.text_frame.paragraphs[0].add_run()
    run0.text = f"SECCIÓN {number}"
    set_run(run0, 14, True, ACCENT)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.5), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, 34, True, WHITE)
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.8))
        run2 = tb2.text_frame.paragraphs[0].add_run()
        run2.text = subtitle
        set_run(run2, 18, False, RGBColor(0xD6, 0xE3, 0xF0))


def content_slide(title, bullets, page, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, 24, True, WHITE)

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(12.2), Inches(5.6))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.space_after = Pt(8)
        run = para.add_run()
        if b.startswith("##"):
            run.text = b[2:].strip()
            set_run(run, 15, True, BLUE)
        elif b.startswith("••"):
            run.text = "    •  " + b[2:].strip()
            set_run(run, 14, False, GRAY)
        else:
            run.text = "•  " + b
            set_run(run, 16, False, DARK)
    if note:
        ntb = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(12), Inches(0.4))
        nr = ntb.text_frame.paragraphs[0].add_run()
        nr.text = note
        set_run(nr, 11, False, GRAY)
    add_footer(slide, page)


def two_col_slide(title, left_title, left_bullets, right_title, right_bullets, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, 24, True, WHITE)

    left = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.25), Inches(6.1), Inches(5.5)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = LIGHT
    left.line.fill.background()
    lt = slide.shapes.add_textbox(Inches(0.65), Inches(1.4), Inches(5.6), Inches(0.45))
    lr = lt.text_frame.paragraphs[0].add_run()
    lr.text = left_title
    set_run(lr, 17, True, BLUE)
    lb = slide.shapes.add_textbox(Inches(0.65), Inches(1.95), Inches(5.6), Inches(4.5))
    ltf = lb.text_frame
    ltf.word_wrap = True
    first = True
    for b in left_bullets:
        para = ltf.paragraphs[0] if first else ltf.add_paragraph()
        first = False
        para.space_after = Pt(7)
        run = para.add_run()
        run.text = "•  " + b
        set_run(run, 13, False, DARK)

    right = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.25), Inches(6.1), Inches(5.5)
    )
    right.fill.solid()
    right.fill.fore_color.rgb = LIGHT
    right.line.fill.background()
    rt = slide.shapes.add_textbox(Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.45))
    rr = rt.text_frame.paragraphs[0].add_run()
    rr.text = right_title
    set_run(rr, 17, True, TEAL)
    rb = slide.shapes.add_textbox(Inches(7.05), Inches(1.95), Inches(5.6), Inches(4.5))
    rtf = rb.text_frame
    rtf.word_wrap = True
    first = True
    for b in right_bullets:
        para = rtf.paragraphs[0] if first else rtf.add_paragraph()
        first = False
        para.space_after = Pt(7)
        run = para.add_run()
        run.text = "•  " + b
        set_run(run, 13, False, DARK)
    add_footer(slide, page)


def cards_slide(title, cards, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, 24, True, WHITE)

    n = len(cards)
    gap = 0.22
    margin = 0.35
    usable = 13.333 - 2 * margin - (n - 1) * gap
    w = usable / n
    colors = [BLUE, TEAL, ACCENT, GREEN, NAVY, ORANGE]
    for i, (ct, lines) in enumerate(cards):
        x = margin + i * (w + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.25), Inches(w), Inches(5.55)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.fill.background()
        top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.25), Inches(w), Inches(0.1)
        )
        top.fill.solid()
        top.fill.fore_color.rgb = colors[i % len(colors)]
        top.line.fill.background()
        ht = slide.shapes.add_textbox(Inches(x + 0.12), Inches(1.45), Inches(w - 0.24), Inches(0.65))
        hr = ht.text_frame.paragraphs[0].add_run()
        hr.text = ct
        set_run(hr, 13, True, colors[i % len(colors)])
        body = slide.shapes.add_textbox(
            Inches(x + 0.12), Inches(2.2), Inches(w - 0.24), Inches(4.3)
        )
        btf = body.text_frame
        btf.word_wrap = True
        first = True
        for line in lines:
            para = btf.paragraphs[0] if first else btf.add_paragraph()
            first = False
            para.space_after = Pt(5)
            run = para.add_run()
            run.text = "• " + line
            set_run(run, 11, False, DARK)
    add_footer(slide, page)


def flow_slide(title, steps, page, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run(run, 24, True, WHITE)

    n = len(steps)
    margin = 0.3
    gap = 0.1
    usable = 13.333 - 2 * margin - (n - 1) * gap
    w = usable / n
    palette = [GREEN, ACCENT, BLUE, ORANGE, RED, RGBColor(0x6C, 0x34, 0x8D)]
    for i, (num, label, desc) in enumerate(steps):
        x = margin + i * (w + gap)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(w), Inches(3.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = palette[i % len(palette)]
        box.line.fill.background()
        nt = slide.shapes.add_textbox(Inches(x + 0.08), Inches(2.2), Inches(w - 0.16), Inches(0.5))
        np = nt.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        nr = np.add_run()
        nr.text = str(num)
        set_run(nr, 26, True, WHITE)
        lt = slide.shapes.add_textbox(Inches(x + 0.08), Inches(2.85), Inches(w - 0.16), Inches(1.0))
        lp = lt.text_frame.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = label
        set_run(lr, 13, True, WHITE)
        dt = slide.shapes.add_textbox(Inches(x + 0.08), Inches(4.0), Inches(w - 0.16), Inches(1.1))
        dtf = dt.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.alignment = PP_ALIGN.CENTER
        dr = dp.add_run()
        dr.text = desc
        set_run(dr, 11, False, WHITE)
    if note:
        ntb = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2))
        ntf = ntb.text_frame
        ntf.word_wrap = True
        nr = ntf.paragraphs[0].add_run()
        nr.text = note
        set_run(nr, 13, False, GRAY)
    add_footer(slide, page)


def agenda_slide(page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Agenda"
    set_run(run, 24, True, WHITE)
    items = [
        ("01", "Visión del módulo y mapa de lo nuevo"),
        ("02", "Flujo operativo del día"),
        ("03", "Jornada y turnos (apertura, parcial, definitivo)"),
        ("04", "Facturación POS (salón, fiscal, cobranza)"),
        ("05", "Informe de gerencia"),
        ("06", "Reportes: descuentos, analítico, ventas, insumos"),
        ("07", "Fórmulas de artículos (costos y stock)"),
        ("08", "Waitry, canjes, viandas, vending, caja/contable"),
        ("09", "Resumen de lo construido"),
    ]
    for i, (num, txt) in enumerate(items):
        y = 1.2 + i * 0.58
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.48), Inches(0.48)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE if i < 6 else TEAL
        circ.line.fill.background()
        nt = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.07), Inches(0.48), Inches(0.35))
        np = nt.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        nr = np.add_run()
        nr.text = num
        set_run(nr, 11, True, WHITE)
        tt = slide.shapes.add_textbox(Inches(1.45), Inches(y + 0.07), Inches(10.5), Inches(0.4))
        tr = tt.text_frame.paragraphs[0].add_run()
        tr.text = txt
        set_run(tr, 16, False, DARK)
    add_footer(slide, page)


def resumen_slide(page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Resumen — qué quedó construido"
    set_run(run, 24, True, WHITE)
    resumen = [
        ("1", "Gobierno diario", "Jornada + turnos por PC + arqueo + PDFs + saneamiento"),
        ("2", "Facturación POS", "Mesas/cuentas, ARCA, cobranza, descuentos, canjes"),
        ("3", "Autoservicio", "Waitry tótem + conciliación Informe Z"),
        ("4", "Marketing / fidelidad", "Canjes premio, tarjeta, VIP y listados"),
        ("5", "Costos y stock", "Fórmulas, opcionales y explosión de insumos"),
        ("6", "Informe gerente", "Dashboard diario: mix, turnos, costos, compras"),
        ("7", "Reportes gestión", "Descuentos, analítico, ventas, insumos, artículos"),
        ("8", "Caja / contable", "Rendiciones, asientos, conciliaciones, auditorías"),
        ("9", "Verticales", "Viandas y máquinas vending"),
        ("10", "Maestros", "PV, mesas, mozos, descuentos, turnos, áreas"),
    ]
    for i, (num, tit, desc) in enumerate(resumen):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 6.45
        y = 1.2 + row * 1.1
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6.2), Inches(0.95)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.fill.background()
        nt = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(0.45), Inches(0.35))
        nr = nt.text_frame.paragraphs[0].add_run()
        nr.text = num
        set_run(nr, 16, True, BLUE)
        tt = slide.shapes.add_textbox(Inches(x + 0.65), Inches(y + 0.12), Inches(5.3), Inches(0.35))
        tr = tt.text_frame.paragraphs[0].add_run()
        tr.text = tit
        set_run(tr, 14, True, DARK)
        dt = slide.shapes.add_textbox(Inches(x + 0.65), Inches(y + 0.48), Inches(5.3), Inches(0.4))
        dr = dt.text_frame.paragraphs[0].add_run()
        dr.text = desc
        set_run(dr, 12, False, GRAY)
    add_footer(slide, page)


# ========== BUILD ==========
title_slide(
    "Módulo Gastronomía",
    "Desarrollo completo · Facturación, turnos, gerencia, descuentos y fórmulas",
    "Anita ERP  ·  Julio 2026",
)

agenda_slide(2)

cards_slide(
    "Mapa de lo nuevo — valor de negocio",
    [
        (
            "Gobierno diario",
            [
                "Jornada por empresa",
                "Turno por terminal",
                "Arqueo y PDFs",
                "Conciliación medios",
                "Saneamiento",
            ],
        ),
        (
            "Facturación",
            [
                "POS mesas/cuentas",
                "Fiscal ARCA",
                "Cobranza multi-medio",
                "Descuentos / NC",
                "Sync Anita",
            ],
        ),
        (
            "Gerencia",
            [
                "Informe gerente",
                "Reporte descuentos",
                "Analítico margen",
                "Ventas / insumos",
                "Facturas del día",
            ],
        ),
        (
            "Costos / stock",
            [
                "Fórmulas / recetas",
                "Opcionales plato",
                "Explosión insumos",
                "Costo vs PVP",
            ],
        ),
    ],
    3,
)

flow_slide(
    "Flujo mínimo operativo del día",
    [
        ("1", "Abrir jornada", "Empresa / día\nde turno"),
        ("2", "Habilitar turno", "Por cada PC\ncon fondo caja"),
        ("3", "Facturar", "Mesas, cuentas,\nWaitry, canjes"),
        ("4", "Cierre parcial", "Arqueo intermedio\n(opcional)"),
        ("5", "Cierre turno", "Definitivo +\nconciliación"),
        ("6", "Cerrar jornada", "Informe Z totem\n+ rendición"),
    ],
    4,
    "Consultas del día: Facturas del día · Cierres · Informe gerente · Descuentos · Artículos vendidos · Saneamiento",
)

# --- TURNOS ---
section_slide("01", "Jornada y turnos", "Apertura, habilitación, parcial, definitivo, central y saneamiento")

content_slide(
    "Tres niveles temporales (concepto clave)",
    [
        "##Jornada — alcance empresa",
        "••Día de turno abierto para toda la empresa (una sola abierta).",
        "••Define la fecha de jornada en los comprobantes.",
        "##Turno operativo — alcance terminal (PC)",
        "••Ventana de caja habilitada en un equipo concreto.",
        "••Cada PC habilita, hace parciales y cierra su turno.",
        "##Cuenta gastronómica — mesa o cuenta libre",
        "••Consumos antes de emitir factura (abierta / cerrada / facturada).",
        "##Fecha de factura",
        "••Siempre el día calendario real; puede diferir de la jornada si el turno cruza medianoche.",
    ],
    6,
)

two_col_slide(
    "Jornada gastronómica",
    "Apertura",
    [
        "Ventas → Gastronomía → Jornada",
        "Empresa + fecha de jornada",
        "Sin jornada abierta no se factura (config)",
        "No abre si quedan turnos sin cerrar",
        "Observación opcional al abrir",
    ],
    "Cierre",
    [
        "Todos los turnos deben estar cerrados",
        "Sin cuentas abiertas con consumos",
        "Conciliación Informe Z Waitry",
        "Proceso contable / rendición Anita",
        "PDF ingresos tótem + historial",
        "Anular último cierre (si no hay rendición)",
    ],
    7,
)

content_slide(
    "Habilitación de turno y cierres",
    [
        "Habilitar en cada PC: turno maestro (mañana/tarde/noche) + monto de caja.",
        "Panel de estado: totales, cuentas pendientes, conciliación por medio / NC / invitaciones.",
        "Cierre parcial: arqueo intermedio con PDF; el turno sigue abierto para vender.",
        "Cierre definitivo: arqueo por medio, NC, invitaciones; PDF + informe por mozo.",
        "Desde el POS: botones de parcial / cerrar turno en la barra del facturador.",
        "Cierre central: supervisión remota de varias terminales.",
        "Histórico de cierres: canjes, tickets tarjeta, corrección de arqueo (permiso).",
        "Saneamiento: diagnóstico, cierre remoto, cuentas pendientes, PDF diagnóstico.",
    ],
    8,
    "Rutas: jornada · habilitacion-turno · cierres-turno · cierre-turno-central · saneamiento-turno",
)

cards_slide(
    "Después del cierre — Caja y Contable",
    [
        (
            "Rendición caja",
            [
                "Rendir turnos/jornadas",
                "Sync Anita rendgastro",
                "caja/rendiciongastronomia",
            ],
        ),
        (
            "Contable",
            [
                "Consulta cierres turno",
                "Conciliación por medio",
                "Diario por punto de venta",
                "Asientos cierre jornada",
            ],
        ),
        (
            "Controles",
            [
                "Conciliación diaria",
                "Auditoría mes vs Anita",
                "Alerta Informe Z faltante",
                "Verificación CAEA",
            ],
        ),
    ],
    9,
)

# --- FACTURACION ---
section_slide("02", "Facturación POS", "Proceso de facturación: mesas, cuentas, fiscal, cobranza y post-venta")

content_slide(
    "Proceso de facturación — pantalla principal del salón",
    [
        "Modo Mesas (plano) o Cuentas libres (sin mesa física).",
        "Apertura de cuenta: cubiertos y mozo (obligatorios según config).",
        "Carga de ítems por SKU o catálogo; opcionales de fórmula por orden.",
        "Cliente de factura / consumidor final / datos manuales según IVA.",
        "Cobranza multi-medio: efectivo, tarjetas, CTG, TOTEM (Waitry), etc.",
        "Atajos: F5 factura (efectivo u otros medios) · F8 facturar con descuento.",
        "Emisión fiscal ARCA: CAE / CAEA con failover y numeración.",
        "Sincronización Anita; ticket / reimpresión desde Facturas del día.",
    ],
    11,
    "Ruta: ventas/gastronomia/proceso-facturacion",
)

two_col_slide(
    "Facturación — reglas y post-venta",
    "Validaciones antes de emitir",
    [
        "Jornada abierta (si es obligatoria)",
        "Turno habilitado en la PC",
        "Config PV + depósitos + lista precios",
        "PV fiscal CAE/CAEA vigente",
        "Cuenta abierta con líneas",
        "Cobranza = total (salvo cortesía)",
        "Canjes premio/fidelidad → solo F8",
        "Sin duplicar órdenes Waitry",
    ],
    "Facturas del día y NC",
    [
        "Consulta de comprobantes de la jornada",
        "Ver canjes asociados (premio/tarjeta)",
        "Reimprimir ticket",
        "Generar nota de crédito",
        "Cambiar medio de pago (permiso)",
        "Base para arqueo y conciliación",
    ],
    12,
)

cards_slide(
    "Configuración y maestros que habilitan la facturación",
    [
        (
            "Punto de venta",
            [
                "Empresa / PC",
                "PV CAE y CAEA",
                "Lista de precios",
                "Tipos de transacción",
                "Depósitos venta/insumos",
                "Waitry habilitado",
            ],
        ),
        (
            "Salón",
            [
                "Mesas + sync Anita",
                "Ubicaciones / sectores",
                "Mozos",
                "Áreas de comanda",
                "Turnos maestros",
            ],
        ),
        (
            "Comercial",
            [
                "Descuentos / invitaciones",
                "Categoría fidelidad",
                "Cliente VIP",
                "Tótems Waitry",
            ],
        ),
        (
            "Fiscal",
            [
                "Pantalla ARCA CAEA",
                "Failover CAE→CAEA",
                "Verificación al cierre",
                "Recuperación comprobantes",
            ],
        ),
    ],
    13,
)

# --- INFORME GERENTE ---
section_slide("03", "Informe de gerencia", "Dashboard con gráficos reales · BIYEMAS · jornada 18/06/2026")

content_slide(
    "Informe gerente — para qué sirve",
    [
        "Foto diaria por empresa y fecha de jornada: qué se vendió, dónde y en qué turno.",
        "Incluye Waitry pagado aún sin facturar cuando la jornada está abierta.",
        "Apoya decisiones de mix, descuentos, costos y compras del día.",
        "Gráficos en pantalla: barras (top artículos día/mes) y tortas (turno, PV, descuentos, recepciones).",
        "Visión gerencial sin entrar al POS.",
        "Permiso: ver-informe-gerente-gastronomia · Ruta: ventas/gastronomia/informe-gerente",
    ],
    15,
)

image_slide(
    "Pantalla del informe gerente",
    IMG / "screenshot_informe.png",
    16,
    caption="Captura de la herramienta en Anita ERP (filtros, totales y tablas Top 10).",
    max_h=5.4,
)

image_slide(
    "KPI de jornada (dato real BIYEMAS 18/06/2026)",
    IMG / "kpi_banner.png",
    17,
    caption="Total neto de la jornada generado por el mismo servicio del informe gerente.",
    max_h=2.8,
)

two_image_slide(
    "Mix de artículos — barras del informe",
    IMG / "articulos_dia.png",
    IMG / "articulos_mes.png",
    18,
    "Top 10 cantidad — día",
    "Top 10 cantidad — mes",
)

two_image_slide(
    "Distribución de ventas — tortas del informe",
    IMG / "turno.png",
    IMG / "puntoventa.png",
    19,
    "Por turno",
    "Por punto de venta",
)

two_image_slide(
    "Descuentos y compras — tortas del informe",
    IMG / "descuento.png",
    IMG / "recepciones_dia.png",
    20,
    "Facturas por código de descuento",
    "Recepciones Anita del día",
)

two_image_slide(
    "Recepciones Anita — día vs mes",
    IMG / "recepciones_dia.png",
    IMG / "recepciones_mes.png",
    21,
    "Día — por proveedor",
    "Mes — por proveedor",
)

cards_slide(
    "Bloques del informe gerente (resumen)",
    [
        (
            "Ventas y mix",
            [
                "Total neto jornada",
                "Estado de jornada",
                "Top 10 por cantidad",
                "Top 10 por valor",
                "Barras día vs mes",
            ],
        ),
        (
            "Distribución",
            [
                "Torta por turno",
                "Torta por PV",
                "Facturas por descuento",
                "Detalle PV",
                "Detalle descuentos",
            ],
        ),
        (
            "Costos",
            [
                "Top 20 precio venta",
                "Costo Anita (stkpre)",
                "Mes anterior vs actual",
            ],
        ),
        (
            "Compras",
            [
                "Recepciones día",
                "Recepciones mes",
                "Por proveedor",
                "Por centro de costo",
            ],
        ),
    ],
    22,
)

# --- REPORTES ---
section_slide(
    "04",
    "Reportes de gestión",
    "Descuentos integrados, analítico, ventas por artículos, insumos y más",
)

content_slide(
    "Reporte de descuentos (integrado)",
    [
        "Análisis de invitaciones y descuentos aplicados en gastronomía.",
        "Cruce por mozos, clientes VIP, clientes internos y códigos de descuento.",
        "Totales y bloques exportables (PDF / Excel / CSV multi-hoja).",
        "Útil para control de cortesías, abuso de descuentos y auditoría comercial.",
        "Cache de consulta para reportes costosos; filtros por empresa y rango de jornada.",
        "Ruta: ventas/gastronomia/descuento-reporte",
        "Complementa la torta de descuentos del Informe gerente (detalle operativo vs dashboard).",
    ],
    23,
)

cards_slide(
    "Suite de reportes gastronomía",
    [
        (
            "Descuentos",
            [
                "Invitaciones / %",
                "Por mozo y VIP",
                "Export multi-bloque",
                "descuento-reporte",
            ],
        ),
        (
            "Analítico",
            [
                "Líneas de venta",
                "Costo y margen",
                "Filtros inteligentes",
                "Export PDF/Excel/CSV",
            ],
        ),
        (
            "Ventas / artículos",
            [
                "Ventas por artículos",
                "Artículos vendidos",
                "Drill a facturas",
                "Drill a mov. stock",
            ],
        ),
        (
            "Insumos",
            [
                "Matriz tipo × día",
                "Consumo proyectado",
                "Ligado a fórmulas",
                "insumos-tipoarticulo",
            ],
        ),
    ],
    25,
)

# --- FORMULAS ---
section_slide("05", "Fórmulas de artículos", "Recetas, costos, opcionales y consumo de insumos al vender")

two_col_slide(
    "ABM de fórmulas (Stock → Fórmula artículo)",
    "Qué se define",
    [
        "Artículo padre / código Anita",
        "Cantidad de porciones",
        "Líneas de insumos + factor costo",
        "Subfórmulas anidadas",
        "Ítems opcionales (AGG/CROWN)",
        "Estados, adjuntos e historia",
        "Export PDF / Excel / CSV",
    ],
    "Integraciones",
    [
        "Sync desde Anita",
        "Vínculo artículos por código",
        "Costo última compra Anita",
        "Artículos compra por insumo",
        "Preview conversión mov. stock",
        "Recálculo costos transferencias",
    ],
    27,
)

content_slide(
    "Cálculo de costos y explosión en operación",
    [
        "##Costo total de la fórmula",
        "••Σ (cantidad × factor × precio última compra), con explosión recursiva de subfórmulas.",
        "••Los opcionales elegidos impactan el costo mostrado en pantalla.",
        "##Consumo al facturar",
        "••Al emitir: baja el plato y los insumos según la fórmula.",
        "••Depósitos de venta e insumos según config del punto de venta.",
        "##Opcionales en el POS",
        "••El cliente elige opcionales por orden; se reflejan en stock y costo.",
        "##Reportes ligados",
        "••Insumos por tipo de artículo · Artículos vendidos · Top 20 costos (informe gerente).",
    ],
    28,
    "stock/formula-articulo + consumo en facturación gastronomía",
)

# --- ECOSISTEMA ---
section_slide("06", "Ecosistema", "Waitry, canjes, viandas, vending y manuales")

cards_slide(
    "Waitry, canjes y verticales",
    [
        (
            "Waitry",
            [
                "ABM tótems",
                "Órdenes en POS",
                "Import por monitor",
                "Cobro TOTEM auto",
                "Informe Z al cierre",
            ],
        ),
        (
            "Canjes Wigos",
            [
                "Premio por cupón (F8)",
                "Fidelidad tarjeta (F8)",
                "Ticket CTG cobranza",
                "Cliente VIP",
                "Facturador marketing",
                "Listado marketing",
            ],
        ),
        (
            "Viandas",
            [
                "Tipos de menú",
                "Usuarios habilitados",
                "Terminal marcha",
                "Reporte / consumos día",
            ],
        ),
        (
            "Vending",
            [
                "ABM máquinas",
                "Rendición ventas",
                "Cierre contable",
                "Sync Anita",
            ],
        ),
    ],
    30,
)

content_slide(
    "Manuales y documentación disponible",
    [
        "Manual de usuario Gastronomía (PDF/Word) — flujo diario, restricciones, Waitry/Wigos.",
        "Manual Canjes Marketing — VIP y facturador de marketing.",
        "Manual Vending — máquinas y rendiciones.",
        "Diagrama ER del módulo (docs/gastronomia).",
        "Ayuda en pantalla desde el menú del módulo.",
    ],
    31,
)

resumen_slide(32)

prs.save(OUT)
print(f"OK: {OUT}")
print(f"Slides: {len(prs.slides)} (esperado {TOTAL})")
